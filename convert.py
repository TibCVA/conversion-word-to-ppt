#!/usr/bin/env python
# -*- coding: utf-8 -*-

import sys
import os
from docx import Document
from docx.text.paragraph import Paragraph
from docx.table import Table
from docx.shared import Pt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor

def px_to_inch(px):
    """Conversion précise des pixels en pouces"""
    return float(px) / 96.0

# Définition explicite des zones avec dimensions minimales garanties
TITLE_ZONE = {
    "x": Inches(px_to_inch(76)),
    "y": Inches(px_to_inch(35)),
    "width": Inches(max(px_to_inch(1382), 1)),  # Largeur minimale de 1 pouce
    "height": Inches(max(px_to_inch(70), 0.5))  # Hauteur minimale de 0.5 pouce
}

SUBTITLE_ZONE = {
    "x": Inches(px_to_inch(76)),
    "y": Inches(px_to_inch(119)),
    "width": Inches(max(px_to_inch(1382), 1)),
    "height": Inches(max(px_to_inch(56), 0.5))
}

CONTENT_ZONE = {
    "x": Inches(px_to_inch(76)),
    "y": Inches(px_to_inch(189)),
    "width": Inches(max(px_to_inch(1382), 1)),
    "height": Inches(max(px_to_inch(425), 1))
}

def validate_template(template_path):
    """Vérifie que le template est valide et contient au moins un layout"""
    try:
        prs = Presentation(template_path)
        if len(prs.slide_layouts) == 0:
            print("Erreur: Le template ne contient aucun layout.")
            return False
            
        # Liste les layouts disponibles
        print("Layouts disponibles dans le template:")
        for i, layout in enumerate(prs.slide_layouts):
            print(f"  - Layout {i}: {layout.name}")
            
        has_blank = any(layout.name == 'Blank' for layout in prs.slide_layouts)
        if not has_blank:
            print("Note: Aucun layout 'Blank' trouvé dans le template. Le premier layout disponible sera utilisé.")
            
        return True
    except Exception as e:
        print(f"Erreur lors de la validation du template: {str(e)}")
        return False

def create_text_frame(shape, text="", font_name="Arial", font_size=11, bold=False, 
                     margin=0.1, auto_size=True):
    """Configuration complète d'une zone de texte"""
    text_frame = shape.text_frame
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE if auto_size else MSO_AUTO_SIZE.NONE
    
    # Définition explicite des marges en pouces
    text_frame.margin_left = text_frame.margin_right = Inches(margin)
    text_frame.margin_top = text_frame.margin_bottom = Inches(margin)
    text_frame.vertical_anchor = MSO_ANCHOR.TOP
    
    if text:
        p = text_frame.paragraphs[0]
        p.text = text
        p.font.name = font_name
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.alignment = PP_ALIGN.LEFT
        
    return text_frame

def add_formatted_text(text_frame, text, font_name="Arial", font_size=11, bold=False, 
                      level=0, bullet=False):
    """Ajoute du texte formaté à un text_frame"""
    p = text_frame.add_paragraph()
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.level = level
    
    if bullet:
        p.bullet.enabled = True
    
    return p

def create_shape_with_text(slide, text, left, top, width, height, font_name="Arial", 
                          font_size=11, bold=False, margin=0.1):
    """Crée une forme avec texte avec des dimensions minimales garanties"""
    width = max(width, Inches(1))  # Largeur minimale
    height = max(height, Inches(0.5))  # Hauteur minimale
    
    shape = slide.shapes.add_textbox(left, top, width, height)
    text_frame = create_text_frame(shape, text, font_name, font_size, bold, margin)
    return shape

def add_table_to_slide(slide, table, left, top, width, height):
    """Ajoute un tableau avec des dimensions minimales garanties"""
    width = max(width, Inches(1))
    height = max(height, Inches(0.5))
    
    rows = len(table.rows)
    cols = len(table.columns)
    
    shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = shape.table
    
    # Copie et formatage du contenu
    for i, row in enumerate(table.rows):
        for j, cell in enumerate(row.cells):
            target_cell = tbl.cell(i, j)
            text = " ".join(p.text for p in cell.paragraphs if p.text)
            
            # Formatage explicite de chaque cellule
            text_frame = target_cell.text_frame
            text_frame.word_wrap = True
            p = text_frame.paragraphs[0]
            p.text = text
            p.font.name = "Arial"
            p.font.size = Pt(10)
            if i == 0:  # En-tête en gras
                p.font.bold = True
    
    return shape

def parse_word_content(doc_path):
    """Parse le contenu Word en préservant tous les styles et la structure"""
    doc = Document(doc_path)
    slides_data = []
    current_slide = None
    current_slide_index = -1
    
    # Fonction helper pour analyser le style d'un paragraphe
    def get_paragraph_style(paragraph):
        style_data = {
            "bullet": False,
            "number": False,
            "level": 0,
            "runs": []
        }
        
        # Détection du niveau et type de liste
        if paragraph._element.pPr is not None:
            if paragraph._element.pPr.numPr is not None:
                ilvl = paragraph._element.pPr.numPr.ilvl
                if ilvl is not None:
                    style_data["level"] = ilvl.val
                numId = paragraph._element.pPr.numPr.numId
                if numId is not None:
                    # Détermine si c'est une puce ou une numérotation
                    # en vérifiant le type de liste dans le document
                    try:
                        num = doc.part.numbering_part.numbering_definitions._numbering.num_list[0]
                        abstract_num_id = num.abstractNumId.val
                        if abstract_num_id == 0:  # Généralement pour les puces
                            style_data["bullet"] = True
                        else:
                            style_data["number"] = True
                    except:
                        style_data["bullet"] = True  # Par défaut, on considère que c'est une puce
        
        # Capture du formatage de chaque segment de texte
        for run in paragraph.runs:
            run_data = {
                "text": run.text,
                "bold": run.bold,
                "italic": run.italic,
                "underline": run.underline,
                "font_size": run.font.size.pt if run.font.size else None
            }
            style_data["runs"].append(run_data)
        
        return style_data
    
    # Parcours séquentiel du document
    elements = []
    for element in doc.element.body:
        if element.tag.endswith('p'):
            elements.append(('paragraph', doc.element.body.find_all(element.tag)[0]))
        elif element.tag.endswith('tbl'):
            elements.append(('table', doc.element.body.find_all(element.tag)[0]))
    
    # Analyse des éléments en préservant l'ordre
    for elem_type, element in elements:
        if elem_type == 'paragraph':
            paragraph = Paragraph(element, doc)
            text = paragraph.text.strip()
            if not text:
                continue
            
            if text.upper().startswith("SLIDE"):
                current_slide_index += 1
                if current_slide is not None:
                    slides_data.append(current_slide)
                current_slide = {
                    "title": "",
                    "subtitle": "",
                    "content": [],
                    "tables": []
                }
            elif current_slide is not None:
                if text.startswith("Titre :"):
                    current_slide["title"] = text[len("Titre :"):].strip()
                elif text.startswith("Sous-titre / Message clé :"):
                    current_slide["subtitle"] = text[len("Sous-titre / Message clé :"):].strip()
                else:
                    style_data = get_paragraph_style(element)
                    para_data = {
                        "text": text,
                        "style": style_data
                    }
                    current_slide["content"].append(para_data)
        
        elif elem_type == 'table' and current_slide is not None:
            current_slide["tables"].append(element)
    
    if current_slide is not None:
        slides_data.append(current_slide)
    
    print(f"Nombre total de slides analysées : {len(slides_data)}")
    return slides_data
    
    if current_slide is not None:
        slides_data.append(current_slide)
    
    # Ajout des tableaux aux slides appropriées
    current_slide_idx = 0
    for table in doc.tables:
        while current_slide_idx < len(slides_data):
            # Vérifie si le tableau appartient à la slide courante
            if table._element.getprevious() is not None:
                prev_text = table._element.getprevious().text.upper()
                if "SLIDE" in prev_text:
                    current_slide_idx += 1
                    continue
            slides_data[current_slide_idx]["tables"].append(table)
            break
    
    return slides_data

def create_slide(prs, slide_data):
    """Crée une slide complète avec tous ses éléments"""
    # Recherche du layout approprié
    layout = None
    for layout_option in prs.slide_layouts:
        if layout_option.name == 'Blank':
            layout = layout_option
            break
    
    if layout is None and len(prs.slide_layouts) > 0:
        layout = prs.slide_layouts[0]
        print(f"Layout 'Blank' non trouvé, utilisation du layout '{layout.name}'")
    
    if layout is None:
        raise ValueError("Aucun layout disponible dans le template PowerPoint")
    
    slide = prs.slides.add_slide(layout)
    
    # Zone de titre
    title_box = create_shape_with_text(
        slide, slide_data["title"],
        TITLE_ZONE["x"], TITLE_ZONE["y"],
        TITLE_ZONE["width"], TITLE_ZONE["height"],
        font_size=22, bold=True
    )
    
    # Zone de sous-titre
    subtitle_box = create_shape_with_text(
        slide, slide_data["subtitle"],
        SUBTITLE_ZONE["x"], SUBTITLE_ZONE["y"],
        SUBTITLE_ZONE["width"], SUBTITLE_ZONE["height"],
        font_size=18
    )
    
    # Position initiale pour le contenu
    current_y = CONTENT_ZONE["y"]
    
    # Ajout du contenu avec formatage préservé
    for content in slide_data["content"]:
        text_box = slide.shapes.add_textbox(
            CONTENT_ZONE["x"], current_y,
            CONTENT_ZONE["width"], Inches(0.3)
        )
        tf = text_box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        style = content["style"]
        
        # Application du formatage des runs
        if style["runs"]:
            p.text = ""  # Clear default text
            for run_data in style["runs"]:
                run = p.add_run()
                run.text = run_data["text"]
                
                # Application des styles du run
                if run_data["bold"]:
                    run.font.bold = True
                if run_data["italic"]:
                    run.font.italic = True
                if run_data["underline"]:
                    run.font.underline = True
                if run_data["font_size"]:
                    run.font.size = Pt(run_data["font_size"])
                run.font.name = "Arial"
        else:
            # Si pas de runs spécifiques, applique le texte directement
            p.text = content["text"]
            p.font.size = Pt(11)
            p.font.name = "Arial"
        
        # Ajout des puces/numéros si nécessaire
        if style["bullet"] or style["number"]:
            p.level = style["level"]
            add_bullet_or_number(p, is_bullet=style["bullet"], level=style["level"])
        
        current_y += Inches(0.3)
    
    # Ajout des tableaux
    for table in slide_data["tables"]:
        if current_y + Inches(1) > (CONTENT_ZONE["y"] + CONTENT_ZONE["height"]):
            print("Warning: Espace insuffisant pour le tableau")
            break
            
        table_height = Inches(len(table.rows) * 0.3)
        shape = add_table_to_slide(
            slide, table,
            CONTENT_ZONE["x"], current_y,
            CONTENT_ZONE["width"], table_height
        )
        current_y += table_height + Inches(0.2)
    
    return slide
    
    # Zone de titre
    title_box = create_shape_with_text(
        slide, slide_data["title"],
        TITLE_ZONE["x"], TITLE_ZONE["y"],
        TITLE_ZONE["width"], TITLE_ZONE["height"],
        font_size=22, bold=True
    )
    
    # Zone de sous-titre
    subtitle_box = create_shape_with_text(
        slide, slide_data["subtitle"],
        SUBTITLE_ZONE["x"], SUBTITLE_ZONE["y"],
        SUBTITLE_ZONE["width"], SUBTITLE_ZONE["height"],
        font_size=18
    )
    
    # Position initiale pour le contenu
    current_y = CONTENT_ZONE["y"]
    
    # Ajout du contenu avec formatage préservé
    for content in slide_data["content"]:
        text_box = slide.shapes.add_textbox(
            CONTENT_ZONE["x"], current_y,
            CONTENT_ZONE["width"], Inches(0.3)
        )
        tf = text_box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        style = content["style"]
        
        # Application du niveau de liste et type de puce
        if style["bullet"] or style["number"]:
            p.level = style["level"]
            if style["bullet"]:
                try:
                    p._pPr.get_or_add_pPr().add_pStyle(val='Bullet')
                    p._pPr.get_or_add_pPr().add_buChar(value='•')
                except:
                    print(f"Impossible d'ajouter une puce à : {content['text'][:30]}...")
            elif style["number"]:
                try:
                    p._pPr.get_or_add_pPr().add_pStyle(val='ListNumber')
                except:
                    print(f"Impossible d'ajouter une numérotation à : {content['text'][:30]}...")
        
        # Application du formatage des runs
        if style["runs"]:
            p.text = ""  # Clear default text
            for run_data in style["runs"]:
                run = p.add_run()
                run.text = run_data["text"]
                
                # Application des styles du run
                if run_data["bold"]:
                    run.font.bold = True
                if run_data["italic"]:
                    run.font.italic = True
                if run_data["underline"]:
                    run.font.underline = True
                if run_data["font_size"]:
                    run.font.size = Pt(run_data["font_size"])
                run.font.name = "Arial"
        else:
            # Si pas de runs spécifiques, applique le texte directement
            p.text = content["text"]
            p.font.size = Pt(11)
            p.font.name = "Arial"
        
        current_y += Inches(0.3)
    
    # Ajout des tableaux
    for table in slide_data["tables"]:
        if current_y + Inches(1) > (CONTENT_ZONE["y"] + CONTENT_ZONE["height"]):
            print("Warning: Espace insuffisant pour le tableau")
            break
            
        table_height = Inches(len(table.rows) * 0.3)
        shape = add_table_to_slide(
            slide, table,
            CONTENT_ZONE["x"], current_y,
            CONTENT_ZONE["width"], table_height
        )
        current_y += table_height + Inches(0.2)
    
    return slide
    
    # Ajout du titre avec dimensions minimales garanties
    title_shape = create_shape_with_text(
        slide, slide_data["title"],
        TITLE_ZONE["x"], TITLE_ZONE["y"],
        TITLE_ZONE["width"], TITLE_ZONE["height"],
        font_size=22, bold=True
    )
    
    # Ajout du sous-titre
    subtitle_shape = create_shape_with_text(
        slide, slide_data["subtitle"],
        SUBTITLE_ZONE["x"], SUBTITLE_ZONE["y"],
        SUBTITLE_ZONE["width"], SUBTITLE_ZONE["height"],
        font_size=18
    )
    
    # Position initiale pour le contenu
    current_y = CONTENT_ZONE["y"]
    content_height = Inches(0.3)  # Hauteur par défaut
    
    # Ajout du contenu
    for content in slide_data["content"]:
        if isinstance(content, dict):  # Paragraphe
            shape = create_shape_with_text(
                slide, content["text"],
                CONTENT_ZONE["x"], current_y,
                CONTENT_ZONE["width"], content_height,
                font_size=11,
                bold=False
            )
            # Ajout des puces si nécessaire
            if content.get("bullet"):
                text_frame = shape.text_frame
                p = text_frame.paragraphs[0]
                try:
                    p._pPr.get_or_add_pPr().add_buNone()  # Réinitialise les puces
                    p._pPr.get_or_add_pPr().add_buChar(value='•')  # Ajoute une puce
                except:
                    print(f"Avertissement: Impossible d'ajouter une puce au paragraphe {content['text'][:30]}...")
            
            current_y += content_height + Inches(0.1)
    
    # Ajout des tableaux
    for table in slide_data["tables"]:
        table_height = Inches(len(table.rows) * 0.3)
        add_table_to_slide(
            slide, table,
            CONTENT_ZONE["x"], current_y,
            CONTENT_ZONE["width"], table_height
        )
        current_y += table_height + Inches(0.2)
    
    return slide

def create_ppt_from_docx(input_docx, template_pptx, output_pptx):
    """Fonction principale de conversion"""
    # Validation des fichiers
    if not os.path.exists(input_docx):
        raise FileNotFoundError(f"Le fichier Word {input_docx} n'existe pas.")
    if not os.path.exists(template_pptx):
        raise FileNotFoundError(f"Le template PowerPoint {template_pptx} n'existe pas.")
    
    # Validation du template
    if not validate_template(template_pptx):
        raise ValueError("Le template PowerPoint n'est pas valide.")
    
    # Extraction du contenu Word
    print("Extraction du contenu Word...")
    slides_data = parse_word_content(input_docx)
    if not slides_data:
        raise ValueError("Aucune slide trouvée dans le document Word.")
    
    # Création de la présentation
    print("Création de la présentation PowerPoint...")
    prs = Presentation(template_pptx)
    
    # Suppression des slides existantes du template
    xml_slides = prs.slides._sldIdLst
    for sld in list(xml_slides):
        xml_slides.remove(sld)
    
    # Création des nouvelles slides
    for i, slide_data in enumerate(slides_data, 1):
        print(f"Création de la slide {i}/{len(slides_data)}...")
        create_slide(prs, slide_data)
    
    # Sauvegarde avec gestion d'erreurs
    try:
        prs.save(output_pptx)
        print(f"Conversion terminée avec succès. {len(slides_data)} slides créées dans {output_pptx}")
    except Exception as e:
        print(f"Erreur lors de la sauvegarde : {str(e)}")
        raise

if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: python convert.py input.docx output.pptx")
        sys.exit(1)
    
    try:
        input_docx = sys.argv[1]
        output_pptx = sys.argv[2]
        template_file = os.path.join(os.path.dirname(os.path.abspath(__file__)), "template_CVA.pptx")
        
        create_ppt_from_docx(input_docx, template_file, output_pptx)
    except Exception as e:
        print(f"Erreur lors de la conversion : {str(e)}")
        sys.exit(1)
